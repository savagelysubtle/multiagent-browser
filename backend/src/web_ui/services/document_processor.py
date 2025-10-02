"""
Enhanced Document Processing Service
Handles creation, conversion, and manipulation of various document formats
"""

import asyncio
import io
import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, BinaryIO
import logging

# Document processing imports
from markitdown import MarkItDown
import magic
import chardet
from docx import Document as DocxDocument
from docx.shared import Inches
import openpyxl
from pptx import Presentation
import PyPDF2
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
import markdown
from markdown_it import MarkdownIt
from bs4 import BeautifulSoup
from jinja2 import Environment, BaseLoader
from fpdf import FPDF
from striprtf.striprtf import rtf_to_text
import json
import csv

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Enhanced document processor with comprehensive format support"""

    def __init__(self):
        self.markitdown = MarkItDown()
        self.md_parser = MarkdownIt()
        self.jinja_env = Environment(loader=BaseLoader())

        # Supported file types and their processors
        self.file_processors = {
            # Text formats
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.markdown': self._process_markdown,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.rtf': self._process_rtf,
            '.csv': self._process_csv,
            '.json': self._process_json,

            # Office formats
            '.docx': self._process_docx,
            '.doc': self._process_doc_legacy,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls_legacy,
            '.pptx': self._process_pptx,
            '.ppt': self._process_ppt_legacy,

            # PDF
            '.pdf': self._process_pdf,

            # Other formats
            '.xml': self._process_xml,
            '.yaml': self._process_yaml,
            '.yml': self._process_yaml,
        }

        # Document templates
        self.templates = {
            'business_letter': self._create_business_letter_template(),
            'report': self._create_report_template(),
            'memo': self._create_memo_template(),
            'invoice': self._create_invoice_template(),
        }

    async def process_file(self, file_path: str, target_format: str = 'markdown') -> Dict[str, Any]:
        """Process any supported file type and convert to target format"""
        try:
            file_path = Path(file_path)
            file_ext = file_path.suffix.lower()

            if file_ext not in self.file_processors:
                # Try using MarkItDown for unknown formats
                return await self._process_with_markitdown(file_path, target_format)

            # Use specific processor for known formats
            processor = self.file_processors[file_ext]
            content_data = await processor(file_path)

            # Convert to target format if needed
            if target_format != 'raw':
                content_data = await self._convert_to_format(content_data, target_format)

            return {
                'success': True,
                'content': content_data,
                'source_format': file_ext[1:],  # Remove dot
                'target_format': target_format,
                'metadata': await self._extract_metadata(file_path)
            }

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'source_format': file_ext[1:] if 'file_ext' in locals() else 'unknown'
            }

    async def create_document_from_template(self, template_name: str, data: Dict[str, Any],
                                          output_format: str = 'pdf') -> Dict[str, Any]:
        """Create a document from a predefined template"""
        try:
            if template_name not in self.templates:
                raise ValueError(f"Template '{template_name}' not found")

            template = self.templates[template_name]

            # Render template with data
            rendered_content = await self._render_template(template, data)

            # Generate document in specified format
            document_data = await self._generate_document(rendered_content, output_format)

            return {
                'success': True,
                'document': document_data,
                'template': template_name,
                'format': output_format,
                'metadata': {
                    'created_at': datetime.now().isoformat(),
                    'template_used': template_name,
                    'data_fields': list(data.keys())
                }
            }

        except Exception as e:
            logger.error(f"Error creating document from template: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    async def convert_document(self, content: str, source_format: str,
                             target_format: str) -> Dict[str, Any]:
        """Convert document content between formats"""
        try:
            # Parse source content
            if source_format == 'markdown':
                parsed_content = self.md_parser.parse(content)
            elif source_format == 'html':
                parsed_content = BeautifulSoup(content, 'html.parser')
            else:
                parsed_content = content

            # Convert to target format
            converted_content = await self._convert_to_format(
                {'content': content, 'parsed': parsed_content},
                target_format
            )

            return {
                'success': True,
                'content': converted_content,
                'source_format': source_format,
                'target_format': target_format
            }

        except Exception as e:
            logger.error(f"Error converting document: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    # File type processors
    async def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text files"""
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'

        with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()

        return {
            'content': content,
            'type': 'text',
            'encoding': encoding
        }

    async def _process_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Process markdown files"""
        text_data = await self._process_text(file_path)
        parsed = self.md_parser.parse(text_data['content'])

        return {
            'content': text_data['content'],
            'parsed': parsed,
            'type': 'markdown',
            'html': markdown.markdown(text_data['content'])
        }

    async def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Process Word documents"""
        doc = DocxDocument(file_path)

        # Extract text content
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = '\n\n'.join(paragraphs)

        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        return {
            'content': content,
            'paragraphs': paragraphs,
            'tables': tables,
            'type': 'docx'
        }

    async def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF files"""
        content = ""
        metadata = {}

        # Try pdfplumber first for better text extraction
        try:
            with pdfplumber.open(file_path) as pdf:
                pages = []
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages.append(page_text)
                        content += page_text + "\n\n"

                metadata = {
                    'page_count': len(pdf.pages),
                    'pages': pages
                }
        except Exception:
            # Fallback to PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                pages = []
                for page in reader.pages:
                    page_text = page.extract_text()
                    pages.append(page_text)
                    content += page_text + "\n\n"

                metadata = {
                    'page_count': len(reader.pages),
                    'pages': pages
                }

        return {
            'content': content.strip(),
            'type': 'pdf',
            'metadata': metadata
        }

    async def _process_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel files"""
        workbook = openpyxl.load_workbook(file_path)
        sheets_data = {}

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []

            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    sheet_data.append(list(row))

            sheets_data[sheet_name] = sheet_data

        # Convert to text representation
        content = ""
        for sheet_name, data in sheets_data.items():
            content += f"Sheet: {sheet_name}\n"
            for row in data:
                content += "\t".join(str(cell) if cell is not None else "" for cell in row) + "\n"
            content += "\n"

        return {
            'content': content,
            'sheets': sheets_data,
            'type': 'xlsx'
        }

    async def _process_with_markitdown(self, file_path: Path, target_format: str) -> Dict[str, Any]:
        """Use MarkItDown for unsupported formats"""
        try:
            result = self.markitdown.convert(str(file_path))
            return {
                'content': result.text_content,
                'type': 'markitdown_converted',
                'original_format': file_path.suffix[1:]
            }
        except Exception as e:
            logger.error(f"MarkItDown conversion failed: {str(e)}")
            raise

    # Document generation methods
    async def _generate_pdf_from_content(self, content: str, metadata: Dict[str, Any] = None) -> bytes:
        """Generate PDF from text content"""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Add title if provided
        if metadata and 'title' in metadata:
            title = Paragraph(metadata['title'], styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))

        # Add content paragraphs
        for paragraph in content.split('\n\n'):
            if paragraph.strip():
                p = Paragraph(paragraph.strip(), styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 6))

        doc.build(story)
        return buffer.getvalue()

    async def _generate_docx_from_content(self, content: str, metadata: Dict[str, Any] = None) -> bytes:
        """Generate DOCX from text content"""
        doc = DocxDocument()

        # Add title if provided
        if metadata and 'title' in metadata:
            title = doc.add_heading(metadata['title'], 0)

        # Add content paragraphs
        for paragraph in content.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    # Template methods
    def _create_business_letter_template(self) -> str:
        return """
{{ date }}

{{ recipient_name }}
{{ recipient_address }}

Dear {{ salutation }},

{{ body }}

Sincerely,

{{ sender_name }}
{{ sender_title }}
{{ sender_company }}
"""

    def _create_report_template(self) -> str:
        return """
# {{ title }}

**Date:** {{ date }}
**Author:** {{ author }}
**Department:** {{ department }}

## Executive Summary

{{ executive_summary }}

## Introduction

{{ introduction }}

## Findings

{{ findings }}

## Recommendations

{{ recommendations }}

## Conclusion

{{ conclusion }}
"""

    async def _render_template(self, template: str, data: Dict[str, Any]) -> str:
        """Render Jinja2 template with data"""
        template_obj = self.jinja_env.from_string(template)
        return template_obj.render(**data)

    async def _convert_to_format(self, content_data: Dict[str, Any], target_format: str) -> Any:
        """Convert content to target format"""
        content = content_data.get('content', '')

        if target_format == 'markdown':
            return content
        elif target_format == 'html':
            if content_data.get('type') == 'markdown':
                return markdown.markdown(content)
            return content
        elif target_format == 'pdf':
            return await self._generate_pdf_from_content(content)
        elif target_format == 'docx':
            return await self._generate_docx_from_content(content)
        else:
            return content

    async def _extract_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract file metadata"""
        stat = file_path.stat()
        return {
            'filename': file_path.name,
            'size': stat.st_size,
            'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'extension': file_path.suffix[1:],
        }

    # Placeholder methods for legacy formats
    async def _process_doc_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .doc files using MarkItDown"""
        return await self._process_with_markitdown(file_path, 'markdown')

    async def _process_xls_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .xls files using MarkItDown"""
        return await self._process_with_markitdown(file_path, 'markdown')

    async def _process_ppt_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .ppt files using MarkItDown"""
        return await self._process_with_markitdown(file_path, 'markdown')

    async def _process_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint files"""
        prs = Presentation(file_path)
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_content.append({
                'slide_number': i + 1,
                'content': '\n'.join(slide_text)
            })

        # Combine all slides into one content string
        content = '\n\n'.join([f"Slide {s['slide_number']}:\n{s['content']}"
                              for s in slides_content if s['content'].strip()])

        return {
            'content': content,
            'slides': slides_content,
            'type': 'pptx'
        }

    async def _process_html(self, file_path: Path) -> Dict[str, Any]:
        """Process HTML files"""
        text_data = await self._process_text(file_path)
        soup = BeautifulSoup(text_data['content'], 'html.parser')

        # Extract text content
        text_content = soup.get_text(separator='\n', strip=True)

        return {
            'content': text_content,
            'html': text_data['content'],
            'parsed': soup,
            'type': 'html'
        }

    async def _process_rtf(self, file_path: Path) -> Dict[str, Any]:
        """Process RTF files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()

        text_content = rtf_to_text(rtf_content)

        return {
            'content': text_content,
            'type': 'rtf'
        }

    async def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """Process CSV files"""
        rows = []
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                rows.append(row)

        # Convert to text representation
        content = '\n'.join(['\t'.join(row) for row in rows])

        return {
            'content': content,
            'rows': rows,
            'type': 'csv'
        }

    async def _process_json(self, file_path: Path) -> Dict[str, Any]:
        """Process JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            json_data = json.load(f)

        # Pretty print JSON as content
        content = json.dumps(json_data, indent=2)

        return {
            'content': content,
            'data': json_data,
            'type': 'json'
        }

    async def _process_xml(self, file_path: Path) -> Dict[str, Any]:
        """Process XML files"""
        text_data = await self._process_text(file_path)
        soup = BeautifulSoup(text_data['content'], 'xml')

        # Extract text content
        text_content = soup.get_text(separator='\n', strip=True)

        return {
            'content': text_content,
            'xml': text_data['content'],
            'parsed': soup,
            'type': 'xml'
        }

    async def _process_yaml(self, file_path: Path) -> Dict[str, Any]:
        """Process YAML files"""
        text_data = await self._process_text(file_path)

        return {
            'content': text_data['content'],
            'type': 'yaml'
        }

    async def _generate_document(self, content: str, output_format: str) -> bytes:
        """Generate document in specified format"""
        if output_format == 'pdf':
            return await self._generate_pdf_from_content(content)
        elif output_format == 'docx':
            return await self._generate_docx_from_content(content)
        elif output_format == 'html':
            html_content = f"<html><body>{markdown.markdown(content)}</body></html>"
            return html_content.encode('utf-8')
        else:
            return content.encode('utf-8')