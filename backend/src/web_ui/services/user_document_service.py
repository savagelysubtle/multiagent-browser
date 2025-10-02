"""
User Document Service - Focused on User Experience
Handles document creation, editing, and format conversion for users
"""

import asyncio
import io
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, BinaryIO
import logging
import json

# Document creation libraries
from docx import Document as DocxDocument
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.styles import Font, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import markdown
from bs4 import BeautifulSoup
from jinja2 import Environment, BaseLoader
import csv
import json

# File processing
import magic
import chardet
from striprtf.striprtf import rtf_to_text
import PyPDF2
import pdfplumber

logger = logging.getLogger(__name__)

class UserDocumentService:
    """Service for user document creation and editing"""

    def __init__(self):
        self.jinja_env = Environment(loader=BaseLoader())

        # Document templates for users
        self.user_templates = {
            'blank_document': {
                'name': 'Blank Document',
                'description': 'Start with a blank document',
                'content': '',
                'format': 'markdown'
            },
            'letter': {
                'name': 'Business Letter',
                'description': 'Professional business letter template',
                'content': self._get_letter_template(),
                'format': 'markdown'
            },
            'resume': {
                'name': 'Resume',
                'description': 'Professional resume template',
                'content': self._get_resume_template(),
                'format': 'markdown'
            },
            'report': {
                'name': 'Report',
                'description': 'Structured report template',
                'content': self._get_report_template(),
                'format': 'markdown'
            },
            'meeting_notes': {
                'name': 'Meeting Notes',
                'description': 'Meeting notes template',
                'content': self._get_meeting_notes_template(),
                'format': 'markdown'
            },
            'project_plan': {
                'name': 'Project Plan',
                'description': 'Project planning document',
                'content': self._get_project_plan_template(),
                'format': 'markdown'
            }
        }

    async def get_available_templates(self) -> List[Dict[str, Any]]:
        """Get list of available document templates for users"""
        templates = []
        for key, template in self.user_templates.items():
            templates.append({
                'id': key,
                'name': template['name'],
                'description': template['description'],
                'format': template['format'],
                'preview': template['content'][:200] + '...' if len(template['content']) > 200 else template['content']
            })
        return templates

    async def create_document_from_template(self, template_id: str, title: str = None) -> Dict[str, Any]:
        """Create a new document from a template"""
        try:
            if template_id not in self.user_templates:
                return {
                    'success': False,
                    'error': f'Template "{template_id}" not found'
                }

            template = self.user_templates[template_id]

            # Replace placeholders with current date and user info
            content = template['content']
            content = content.replace('{{date}}', datetime.now().strftime('%B %d, %Y'))
            content = content.replace('{{title}}', title or 'Untitled Document')

            return {
                'success': True,
                'document': {
                    'title': title or template['name'],
                    'content': content,
                    'format': template['format'],
                    'template_used': template_id,
                    'created_at': datetime.now().isoformat()
                }
            }

        except Exception as e:
            logger.error(f"Error creating document from template: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    async def export_document(self, content: str, format: str, title: str = "document") -> Dict[str, Any]:
        """Export document to various formats"""
        try:
            if format == 'pdf':
                document_bytes = await self._create_pdf(content, title)
                return {
                    'success': True,
                    'data': document_bytes,
                    'mime_type': 'application/pdf',
                    'filename': f'{title}.pdf'
                }

            elif format == 'docx':
                document_bytes = await self._create_docx(content, title)
                return {
                    'success': True,
                    'data': document_bytes,
                    'mime_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'filename': f'{title}.docx'
                }

            elif format == 'html':
                html_content = markdown.markdown(content)
                html_doc = f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <title>{title}</title>
                    <meta charset="utf-8">
                    <style>
                        body {{ font-family: Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }}
                        h1, h2, h3 {{ color: #333; }}
                        p {{ line-height: 1.6; }}
                    </style>
                </head>
                <body>
                    {html_content}
                </body>
                </html>
                """
                return {
                    'success': True,
                    'data': html_doc.encode('utf-8'),
                    'mime_type': 'text/html',
                    'filename': f'{title}.html'
                }

            elif format == 'txt':
                # Convert markdown to plain text
                soup = BeautifulSoup(markdown.markdown(content), 'html.parser')
                plain_text = soup.get_text()
                return {
                    'success': True,
                    'data': plain_text.encode('utf-8'),
                    'mime_type': 'text/plain',
                    'filename': f'{title}.txt'
                }

            else:
                return {
                    'success': False,
                    'error': f'Unsupported export format: {format}'
                }

        except Exception as e:
            logger.error(f"Error exporting document: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    async def import_document(self, file_path: str) -> Dict[str, Any]:
        """Import document from file and convert to editable format"""
        try:
            file_path = Path(file_path)
            file_ext = file_path.suffix.lower()

            if file_ext == '.txt':
                content = await self._import_text(file_path)
            elif file_ext == '.md' or file_ext == '.markdown':
                content = await self._import_markdown(file_path)
            elif file_ext == '.docx':
                content = await self._import_docx(file_path)
            elif file_ext == '.pdf':
                content = await self._import_pdf(file_path)
            elif file_ext == '.html' or file_ext == '.htm':
                content = await self._import_html(file_path)
            elif file_ext == '.rtf':
                content = await self._import_rtf(file_path)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file format: {file_ext}'
                }

            return {
                'success': True,
                'document': {
                    'title': file_path.stem,
                    'content': content,
                    'format': 'markdown',
                    'original_format': file_ext[1:],
                    'imported_at': datetime.now().isoformat()
                }
            }

        except Exception as e:
            logger.error(f"Error importing document: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    async def create_spreadsheet(self, data: List[List[str]], title: str = "spreadsheet") -> Dict[str, Any]:
        """Create Excel spreadsheet from data"""
        try:
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            sheet.title = "Sheet1"

            # Add data to sheet
            for row_idx, row_data in enumerate(data, 1):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = sheet.cell(row=row_idx, column=col_idx, value=cell_value)

                    # Style header row
                    if row_idx == 1:
                        cell.font = Font(bold=True)
                        cell.alignment = Alignment(horizontal='center')

            # Auto-adjust column widths
            for column in sheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                sheet.column_dimensions[column_letter].width = adjusted_width

            # Save to bytes
            buffer = io.BytesIO()
            workbook.save(buffer)

            return {
                'success': True,
                'data': buffer.getvalue(),
                'mime_type': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'filename': f'{title}.xlsx'
            }

        except Exception as e:
            logger.error(f"Error creating spreadsheet: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    async def create_presentation(self, slides_data: List[Dict[str, str]], title: str = "presentation") -> Dict[str, Any]:
        """Create PowerPoint presentation from slide data"""
        try:
            prs = Presentation()

            for slide_data in slides_data:
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)

                # Set title
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get('title', 'Slide Title')

                # Set content
                content_shape = slide.placeholders[1]
                content_shape.text = slide_data.get('content', '')

            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)

            return {
                'success': True,
                'data': buffer.getvalue(),
                'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'filename': f'{title}.pptx'
            }

        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    # Document creation methods
    async def _create_pdf(self, content: str, title: str) -> bytes:
        """Create PDF from markdown content"""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=1*inch)
        styles = getSampleStyleSheet()
        story = []

        # Add title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            alignment=1,  # Center alignment
            spaceAfter=30
        )
        story.append(Paragraph(title, title_style))

        # Convert markdown to HTML then to paragraphs
        html_content = markdown.markdown(content)
        soup = BeautifulSoup(html_content, 'html.parser')

        for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol']):
            if element.name.startswith('h'):
                # Handle headers
                level = int(element.name[1])
                style = styles[f'Heading{min(level, 4)}']
                story.append(Paragraph(element.get_text(), style))
            elif element.name == 'p':
                # Handle paragraphs
                story.append(Paragraph(element.get_text(), styles['Normal']))
            elif element.name in ['ul', 'ol']:
                # Handle lists
                for li in element.find_all('li'):
                    bullet_style = ParagraphStyle(
                        'BulletStyle',
                        parent=styles['Normal'],
                        leftIndent=20,
                        bulletIndent=10
                    )
                    story.append(Paragraph(f"• {li.get_text()}", bullet_style))

            story.append(Spacer(1, 6))

        doc.build(story)
        return buffer.getvalue()

    async def _create_docx(self, content: str, title: str) -> bytes:
        """Create DOCX from markdown content"""
        doc = DocxDocument()

        # Add title
        title_paragraph = doc.add_heading(title, 0)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Convert markdown to HTML then process
        html_content = markdown.markdown(content)
        soup = BeautifulSoup(html_content, 'html.parser')

        for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol']):
            if element.name.startswith('h'):
                # Handle headers
                level = int(element.name[1])
                doc.add_heading(element.get_text(), level)
            elif element.name == 'p':
                # Handle paragraphs
                doc.add_paragraph(element.get_text())
            elif element.name in ['ul', 'ol']:
                # Handle lists
                for li in element.find_all('li'):
                    doc.add_paragraph(li.get_text(), style='List Bullet')

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    # Import methods
    async def _import_text(self, file_path: Path) -> str:
        """Import plain text file"""
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'

        with open(file_path, 'r', encoding=encoding) as f:
            return f.read()

    async def _import_markdown(self, file_path: Path) -> str:
        """Import markdown file"""
        return await self._import_text(file_path)

    async def _import_docx(self, file_path: Path) -> str:
        """Import DOCX file and convert to markdown"""
        doc = DocxDocument(file_path)
        markdown_content = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Check if it's a heading (simple heuristic)
                if paragraph.style.name.startswith('Heading'):
                    level = 1  # Default level
                    try:
                        level = int(paragraph.style.name.split()[-1])
                    except:
                        pass
                    markdown_content.append(f"{'#' * level} {text}")
                else:
                    markdown_content.append(text)
                markdown_content.append("")  # Add blank line

        return "\n".join(markdown_content)

    async def _import_pdf(self, file_path: Path) -> str:
        """Import PDF file and extract text"""
        content = ""

        try:
            # Try pdfplumber first
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        content += page_text + "\n\n"
        except Exception:
            # Fallback to PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text() + "\n\n"

        return content.strip()

    async def _import_html(self, file_path: Path) -> str:
        """Import HTML file and convert to markdown"""
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Convert to markdown-like format
        markdown_content = []

        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol']):
            if element.name.startswith('h'):
                level = int(element.name[1])
                markdown_content.append(f"{'#' * level} {element.get_text().strip()}")
            elif element.name == 'p':
                text = element.get_text().strip()
                if text:
                    markdown_content.append(text)
            elif element.name in ['ul', 'ol']:
                for li in element.find_all('li'):
                    markdown_content.append(f"- {li.get_text().strip()}")

            markdown_content.append("")  # Add blank line

        return "\n".join(markdown_content)

    async def _import_rtf(self, file_path: Path) -> str:
        """Import RTF file and convert to plain text"""
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()

        return rtf_to_text(rtf_content)

    # Placeholder methods for legacy formats (removed textract dependency)
    async def _process_doc_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .doc files - requires conversion to .docx first"""
        return {
            'success': False,
            'error': 'Legacy .doc files not supported. Please convert to .docx format first.',
            'suggestion': 'Use Microsoft Word or LibreOffice to save as .docx format'
        }

    async def _process_xls_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .xls files - requires conversion to .xlsx first"""
        return {
            'success': False,
            'error': 'Legacy .xls files not supported. Please convert to .xlsx format first.',
            'suggestion': 'Use Microsoft Excel or LibreOffice Calc to save as .xlsx format'
        }

    async def _process_ppt_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Process legacy .ppt files - requires conversion to .pptx first"""
        return {
            'success': False,
            'error': 'Legacy .ppt files not supported. Please convert to .pptx format first.',
            'suggestion': 'Use Microsoft PowerPoint or LibreOffice Impress to save as .pptx format'
        }

    # Template methods
    def _get_letter_template(self) -> str:
        return """{{date}}

[Recipient Name]
[Recipient Address]
[City, State ZIP Code]

Dear [Recipient Name],

[Write your letter content here. Start with the purpose of your letter, provide any necessary details in the body paragraphs, and end with a clear call to action or next steps.]

[Second paragraph if needed for additional information or supporting details.]

Thank you for your time and consideration.

Sincerely,

[Your Name]
[Your Title]
[Your Company]
[Your Contact Information]"""

    def _get_resume_template(self) -> str:
        return """# [Your Full Name]

**Email:** [your.email@example.com] | **Phone:** [Your Phone Number]
**LinkedIn:** [Your LinkedIn Profile] | **Location:** [City, State]

---

## Professional Summary

[2-3 sentences describing your professional background, key skills, and career objectives]

---

## Experience

### [Job Title] - [Company Name]
*[Start Date] - [End Date]*

- [Key achievement or responsibility]
- [Another achievement with quantifiable results]
- [Third key responsibility or accomplishment]

### [Previous Job Title] - [Previous Company]
*[Start Date] - [End Date]*

- [Key achievement or responsibility]
- [Another achievement with quantifiable results]

---

## Education

### [Degree] in [Field of Study]
**[University Name]** - [Graduation Year]

---

## Skills

**Technical Skills:** [List relevant technical skills]
**Languages:** [Programming languages, spoken languages, etc.]
**Tools & Platforms:** [Software, platforms, frameworks]

---

## Projects

### [Project Name]
[Brief description of the project and your role]

### [Another Project]
[Brief description and impact]"""

    def _get_report_template(self) -> str:
        return """# {{title}}

**Date:** {{date}}
**Prepared by:** [Your Name]
**Department:** [Your Department]

---

## Executive Summary

[Provide a brief overview of the report's key findings and recommendations]

## Introduction

### Purpose
[Explain the purpose and scope of this report]

### Background
[Provide context and background information]

## Methodology

[Describe how data was collected and analyzed]

## Findings

### Key Finding 1
[Detail your first major finding]

### Key Finding 2
[Detail your second major finding]

### Key Finding 3
[Detail your third major finding]

## Analysis

[Provide detailed analysis of your findings]

## Recommendations

1. **Recommendation 1:** [Specific actionable recommendation]
2. **Recommendation 2:** [Another specific recommendation]
3. **Recommendation 3:** [Third recommendation]

## Conclusion

[Summarize the key points and next steps]

---

## Appendices

### Appendix A: Data Sources
[List your data sources]

### Appendix B: Additional Information
[Any supporting information]"""

    def _get_meeting_notes_template(self) -> str:
        return """# Meeting Notes: [Meeting Title]

**Date:** {{date}}
**Time:** [Meeting Time]
**Location:** [Meeting Location/Platform]
**Attendees:** [List of attendees]

---

## Agenda

1. [Agenda item 1]
2. [Agenda item 2]
3. [Agenda item 3]

---

## Discussion Points

### [Topic 1]
- [Key discussion point]
- [Decision made or action needed]

### [Topic 2]
- [Key discussion point]
- [Decision made or action needed]

### [Topic 3]
- [Key discussion point]
- [Decision made or action needed]

---

## Action Items

| Action Item | Assigned To | Due Date | Status |
|-------------|-------------|----------|---------|
| [Action 1] | [Person] | [Date] | [Status] |
| [Action 2] | [Person] | [Date] | [Status] |
| [Action 3] | [Person] | [Date] | [Status] |

---

## Next Steps

- [Next step 1]
- [Next step 2]
- [Next step 3]

**Next Meeting:** [Date and time of next meeting]"""

    def _get_project_plan_template(self) -> str:
        return """# Project Plan: {{title}}

**Project Manager:** [Your Name]
**Start Date:** [Project Start Date]
**End Date:** [Project End Date]
**Last Updated:** {{date}}

---

## Project Overview

### Objective
[Define the main objective of the project]

### Scope
[Describe what is included and excluded from the project]

### Success Criteria
- [Success criterion 1]
- [Success criterion 2]
- [Success criterion 3]

---

## Stakeholders

| Name | Role | Responsibilities | Contact |
|------|------|------------------|---------|
| [Name] | [Role] | [Responsibilities] | [Email] |
| [Name] | [Role] | [Responsibilities] | [Email] |

---

## Project Timeline

### Phase 1: [Phase Name]
**Duration:** [Start Date] - [End Date]

- [ ] [Task 1]
- [ ] [Task 2]
- [ ] [Task 3]

### Phase 2: [Phase Name]
**Duration:** [Start Date] - [End Date]

- [ ] [Task 1]
- [ ] [Task 2]
- [ ] [Task 3]

### Phase 3: [Phase Name]
**Duration:** [Start Date] - [End Date]

- [ ] [Task 1]
- [ ] [Task 2]
- [ ] [Task 3]

---

## Resources Required

### Human Resources
- [Role/Person and allocation]
- [Role/Person and allocation]

### Technical Resources
- [Tool/Software needed]
- [Equipment needed]

### Budget
- **Total Budget:** [Amount]
- **Budget Breakdown:**
  - [Category]: [Amount]
  - [Category]: [Amount]

---

## Risk Management

| Risk | Probability | Impact | Mitigation Strategy |
|------|-------------|---------|-------------------|
| [Risk 1] | [High/Med/Low] | [High/Med/Low] | [Strategy] |
| [Risk 2] | [High/Med/Low] | [High/Med/Low] | [Strategy] |

---

## Communication Plan

- **Status Updates:** [Frequency and method]
- **Team Meetings:** [Schedule]
- **Stakeholder Reports:** [Frequency and format]"""